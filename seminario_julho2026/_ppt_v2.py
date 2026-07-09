# -*- coding: utf-8 -*-
# ==============================================================================
#  _ppt_v2.py — Gera a versao revisada do seminario (Apresentacao..._v2.pptx)
#  a partir da apresentacao atual, aplicando as melhorias pedidas (contextualizacao,
#  motivacao, conceitos, RSL, lacunas, arquitetura por etapas, fusao, pseudocodigo,
#  split/AUC/ablacao, baseline e discussao). SO usa dados REAIS da pesquisa.
#  Nao sobrescreve o original.
# ==============================================================================
import copy
import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import _ppt as H   # toolkit existente (duplicate_slide, fill_definitions, etc.)

VINHO = RGBColor(0x81, 0x05, 0x39)   # cor do titulo (padrao do deck)
CINZA = RGBColor(0x60, 0x60, 0x60)
CORPO = RGBColor(0x30, 0x30, 0x30)

SRC = "Apresentacao_Seminario_Julho2026.pptx"
OUT = "Apresentacao_Seminario_Julho2026_v2.pptx"

prs = Presentation(SRC)


# ── helpers de posicionamento ────────────────────────────────────────────────
def idx(sub):
    return H.slide_index_by_text(prs, sub)

def move_after(moving_idx, anchor_sub):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[moving_idx]
    lst.remove(el)
    a = H.slide_index_by_text(prs, anchor_sub)   # recomputa apos remocao
    anchor_el = list(lst)[a]
    pos = list(lst).index(anchor_el) + 1
    lst.insert(pos, el)

def fmt_titulo(slide, size=32):
    """Padroniza o titulo (Arial Black, vinho, tamanho fixo) e DESLIGA o autofit
    para o texto nao encolher — causa da 'formatacao quebrada' nos slides novos."""
    t = H.find_title(slide)
    if t is None:
        return
    tf = t.text_frame
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.word_wrap = True
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = "Arial Black"; r.font.bold = True
            r.font.size = Pt(size); r.font.color.rgb = VINHO

def normaliza_corpo(tb, cor=CORPO):
    tf = tb.text_frame
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.word_wrap = True

def novo_conteudo(titulo, blocos, depois_de, tsize=32, **kw):
    """Cria um slide de conteudo (template 'Perguntas de Pesquisa'), preenche,
    padroniza o titulo e posiciona logo apos o slide `depois_de`."""
    slide, tb = H.add_content_slide(prs, titulo, blocos, **kw)
    fmt_titulo(slide, tsize)
    normaliza_corpo(tb)
    novo_idx = list(prs.slides).index(slide)
    move_after(novo_idx, depois_de)
    return slide

def add_algoritmo(titulo, require, ensure, linhas, depois_de):
    """Slide de pseudocodigo no estilo 'Algorithm' (Require/Ensure, linhas
    numeradas, comentarios com ▷), em fonte monoespacada."""
    ref = H.slide_index_by_text(prs, "Perguntas de Pesquisa")
    slide = H.duplicate_slide(prs, ref)
    title = H.clear_body(slide)
    H.set_title_text(title, titulo)
    fmt_titulo(slide, 30)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.95), Inches(5.8))
    tf = tb.text_frame; tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    MONO, SZ = "Consolas", 12.5
    primeiro = [True]
    def linha(runs, before=0, after=1):
        p = tf.paragraphs[0] if primeiro[0] else tf.add_paragraph()
        primeiro[0] = False
        p.space_before = Pt(before); p.space_after = Pt(after)
        pPr = p._p.get_or_add_pPr()
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))
        for (txt, bold, ital, color) in runs:
            r = p.add_run(); r.text = txt
            r.font.name = MONO; r.font.size = Pt(SZ)
            r.font.bold = bold; r.font.italic = ital
            if color is not None:
                r.font.color.rgb = color
        return p
    linha([("Require: ", True, False, VINHO), (require, False, False, CORPO)])
    linha([("Ensure:  ", True, False, VINHO), (ensure, False, False, CORPO)], after=6)
    for (num, indent, code, comment) in linhas:
        runs = [((f"{num:>2}: " if num else "    "), False, False, CINZA),
                (("   " * indent) + code, False, False, CORPO)]
        if comment:
            runs.append(("   ▷ " + comment, False, True, CINZA))
        linha(runs)
    move_after(list(prs.slides).index(slide), depois_de)
    return slide

def reformatar(sub, titulo, blocos, **kw):
    """Limpa o corpo de um slide existente e reescreve com definicoes limpas."""
    i = idx(sub)
    slide = list(prs.slides)[i]
    title = H.clear_body(slide)
    H.set_title_text(title, titulo)
    tb = slide.shapes.add_textbox(Inches(kw.pop("left", 0.9)), Inches(kw.pop("top", 1.7)),
                                  Inches(kw.pop("width", 11.4)), Inches(kw.pop("height", 5.2)))
    H.fill_definitions(tb.text_frame, blocos, **kw)
    return slide


GRAF = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "Mestrado_PETR4")

def add_imagem(titulo, arquivo, depois_de, legenda="", blocos_topo=None):
    """Slide com imagem (grafico real) centralizada; titulo padronizado e,
    opcionalmente, um bloco de texto explicativo acima."""
    ref = H.slide_index_by_text(prs, "Perguntas de Pesquisa")
    slide = H.duplicate_slide(prs, ref)
    title = H.clear_body(slide)
    H.set_title_text(title, titulo); fmt_titulo(slide, 32)
    top = 1.75
    if blocos_topo:
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.5), Inches(2.0))
        H.fill_definitions(tb.text_frame, blocos_topo, sz_label=16, sz_desc=14, space=7)
        normaliza_corpo(tb)
        top = 4.0
    h_in = min(4.4, 6.95 - top - (0.35 if legenda else 0.0))
    pic = slide.shapes.add_picture(os.path.join(GRAF, arquivo), Inches(0), Inches(top), height=Inches(h_in))
    pic.left = int((prs.slide_width - pic.width) / 2)
    if legenda:
        cap = slide.shapes.add_textbox(Inches(0.9), Inches(top + h_in + 0.03), Inches(11.5), Inches(0.4))
        r = cap.text_frame.paragraphs[0].add_run(); r.text = legenda
        r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = CINZA
    move_after(list(prs.slides).index(slide), depois_de)
    return slide

def add_tabela(titulo, nota, cabecalho, linhas, depois_de, larguras=None, destaque=1):
    """Slide com uma tabela nativa (numeros reais)."""
    ref = H.slide_index_by_text(prs, "Perguntas de Pesquisa")
    slide = H.duplicate_slide(prs, ref)
    title = H.clear_body(slide)
    H.set_title_text(title, titulo); fmt_titulo(slide, 32)
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.95))
    H.fill_definitions(tb.text_frame, [("Como ler", nota)], sz_label=15, sz_desc=14, space=4)
    normaliza_corpo(tb)
    nr, nc = len(linhas) + 1, len(cabecalho)
    gt = slide.shapes.add_table(nr, nc, Inches(1.4), Inches(2.65), Inches(10.5), Inches(3.9)).table
    if larguras:
        for j, wv in enumerate(larguras):
            gt.columns[j].width = Inches(wv)
    for j, h in enumerate(cabecalho):
        c = gt.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True; r.font.size = Pt(14)
    for i, linha in enumerate(linhas, start=1):
        for j, v in enumerate(linha):
            c = gt.cell(i, j); c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(13); r.font.bold = (i == destaque)
    move_after(list(prs.slides).index(slide), depois_de)
    return slide


# ============================================================================
#  1) INTRODUCAO — contextualizacao + motivacao (inserir apos o divisor)
# ============================================================================
novo_conteudo(
    "Contextualização — o problema",
    [("Um problema central em finanças quantitativas.",
      "Prever a direção e o risco de um ativo orienta investimento, gestão de risco e decisão."),
     ("Mercados são dinâmicos, voláteis e movidos por informação.",
      "O preço reflete preços passados, mas também notícias e fatores macroeconômicos."),
     ("As notícias podem anteceder o movimento do preço.",
      "Parte dessa informação não está nos números — está no texto, antes de virar cotação."),
     ("Previsão é, por natureza, uma tarefa difícil.",
      "Baixa relação sinal-ruído e eventos extremos (choques) limitam qualquer modelo.")],
    depois_de="Introdução", sz_label=17, sz_desc=15, space=13)

novo_conteudo(
    "Contextualização — por que a PETR4?",
    [("A ação mais negociada e mais noticiada da B3.",
      "Alta liquidez e cobertura de mídia constante — ideal para estudar sentimento."),
     ("Duplamente sensível a notícias.",
      "Sofre com a empresa (governança, dividendos, política de preços) E com o petróleo (Brent, geopolítica)."),
     ("Ativo altamente politizado.",
      "Trocas de comando e intervenções tornam o texto especialmente informativo sobre o risco."),
     ("Foco no ativo, não no índice.",
      "A literatura costuma prever índices inteiros; aqui, o risco idiossincrático de UM ativo.")],
    depois_de="Contextualização — o problema", sz_label=17, sz_desc=15, space=13)

novo_conteudo(
    "Motivação — oportunidade e lacuna",
    [("Avanços recentes em IA de linguagem.",
      "Transformers (BERT/FinBERT) leem o contexto do texto — muito além de contar palavras."),
     ("Risco tem estrutura própria e modelável.",
      "O GARCH captura o agrupamento de volatilidade (dias turbulentos vêm em sequência)."),
     ("A lacuna: unir os dois mundos, em português, para um ativo.",
      "Poucos estudos fundem NLP de ponta com risco econométrico (GARCH) aplicados à PETR4."),
     ("A aposta desta pesquisa.",
      "Fundir sentimento (FinBERT-PT-BR) + risco (GARCH) + preços e medir se isso ajuda a prever direção e volatilidade.")],
    depois_de="Choque Informacional", sz_label=17, sz_desc=15, space=13)


# ============================================================================
#  3) CONCEITOS FUNDAMENTAIS (slide 11) — reformatar limpo
# ============================================================================
reformatar("Conceitos Fundamentais", "Conceitos Fundamentais",
    [("FinBERT-PT-BR", "Transformer (BERTimbau ajustado a finanças) que lê a manchete e devolve um sentimento contextual em [−1, +1]."),
     ("GARCH(1,1)", "Modelo econométrico que estima a volatilidade condicional (o risco) e captura o agrupamento de turbulência."),
     ("XGBoost", "Conjunto de árvores de decisão (gradient boosting); é o classificador PRINCIPAL da direção (alta/baixa)."),
     ("SVM", "Classificador que separa as classes maximizando a margem; usado como COMPARAÇÃO."),
     ("Data Fusion", "Une num único vetor os atributos heterogêneos — sentimento, risco (GARCH), retorno e categorias."),
     ("Walk-Forward Validation", "Validação cronológica em janelas deslizantes: treina no passado, testa no futuro, sem vazamento.")],
    sz_label=18, sz_desc=15, space=12)


# ============================================================================
#  6) PROTOCOLO DA RSL (slide 13) — reformatar limpo (numeros reconciliados)
# ============================================================================
reformatar("Protocolo da Revisão", "Protocolo da Revisão Sistemática (RSL)",
    [("Bases consultadas", "ACM Digital Library, IEEE Xplore, Scopus, Web of Science e Periódicos CAPES."),
     ("Termos de busca (PT/EN)", "“previsão de preço de ações”, “volatilidade”, “sentimento de notícias”; “stock price prediction”, “news sentiment”, “stock volatility forecasting”, “financial text mining”."),
     ("Janela da LITERATURA", "2007–2026 — do estudo seminal de Tetlock (2007) aos trabalhos mais recentes."),
     ("Janela dos DADOS empíricos", "2016–2026 — coleta das notícias e das cotações da PETR4."),
     ("Critérios de exclusão", "sem revisão por pares; sem aplicação a mercados reais; sem uso de PLN."),
     ("Funil (PRISMA)", "452 publicações identificadas → 423 excluídas → 29 estudos selecionados.")],
    sz_label=17, sz_desc=14.5, space=12)


# ============================================================================
#  4) DIMENSOES DA COMPARACAO (apos a tabela de trabalhos relacionados)
# ============================================================================
novo_conteudo(
    "Como ler a tabela comparativa",
    [("Ativo × Índice", "A maioria prevê índices (IBOVESPA, S&P). Aqui: um ATIVO específico (PETR4)."),
     ("Texto: dicionário × Transformer", "Trabalhos antigos contam palavras (TF-IDF/léxico). Aqui: FinBERT-PT-BR (contexto e ironia)."),
     ("Idioma", "Predomínio do inglês; só 5 estudos em português. Esta pesquisa é em PT-BR."),
     ("Risco explícito (GARCH)", "Poucos modelam a volatilidade condicional junto do texto. Aqui, o GARCH entra na fusão."),
     ("Alvo: direção × volatilidade", "Muitos preveem só a direção. Aqui, direção E volatilidade (o risco)."),
     ("Fusão de dados", "Coluna a coluna, a diferença é a UNIÃO: NLP de ponta + risco econométrico + preços, no mesmo vetor.")],
    depois_de="síntese das 25 publicações", sz_label=16.5, sz_desc=14.5, space=11)


# ============================================================================
#  5) LACUNAS (slides 16-18) — reescrever texto mais claro (mantendo o reveal)
# ============================================================================
GAP = {
 "GRANULARIDADE (ATIVO)":
   "Sair dos índices inteiros (ex.: IBOVESPA) e mirar o risco idiossincrático de um ativo altamente politizado — a PETR4.",
 "SEMÂNTICA (NLP)":
   "Substituir dicionários estáticos e contagem de palavras (TF-IDF) por Transformers de ponta (FinBERT-PT-BR), que entendem contexto, ironia e jargão do português.",
 "ARQUITETURA (DATA FUSION)":
   "Unificar, numa só arquitetura, a IA textual e a predição não linear (XGBoost/SVM) com o risco econométrico (GARCH) — o que a literatura ainda trata de forma fragmentada.",
}
for s in list(prs.slides):
    for sh in list(s.shapes):
        if sh.has_text_frame:
            for rot, novo in GAP.items():
                # o rotulo (ex.: 'GRANULARIDADE (ATIVO)') fica num shape; a descricao noutro
                pass
# Reescreve as descricoes das lacunas (procura o shape cujo texto comeca com a frase antiga)
_ANTIGAS = [
 "Transição das avaliações genéricas de índices",
 "Superação de dicionários estáticos",
 "Unificação integrando Inteligência Artificial textual",
]
_NOVAS = [GAP["GRANULARIDADE (ATIVO)"], GAP["SEMÂNTICA (NLP)"], GAP["ARQUITETURA (DATA FUSION)"]]
for s in list(prs.slides):
    for sh in list(s.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            t = sh.text_frame.text.strip()
            for a, n in zip(_ANTIGAS, _NOVAS):
                if t.startswith(a[:30]):
                    tf = sh.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    r = p.add_run(); r.text = n; r.font.size = Pt(16)


# ============================================================================
#  6) METODOLOGIA (slide 20) — reformatar limpo
# ============================================================================
reformatar("Metodologia da Pesquisa", "Metodologia da Pesquisa",
    [("Método", "Experimental — compara, sob condições controladas, modelos COM e SEM a informação de sentimento."),
     ("Finalidade", "Aplicada — constrói um artefato computacional que produz previsões verificáveis."),
     ("Unidade de análise", "O pregão (dia de negociação): as notícias do dia viram um índice diário → uma previsão por pregão."),
     ("Pressuposto central", "Causalidade temporal — só a informação disponível ANTES do pregão pode prevê-lo (atributos defasados em t−1)."),
     ("Alvo duplo", "Direção do fechamento (alta/baixa) E volatilidade (o risco).")],
    sz_label=18, sz_desc=15, space=13)


# ============================================================================
#  7) ARQUITETURA POR ETAPAS (apos o slide 23 da arquitetura)
# ============================================================================
novo_conteudo(
    "Arquitetura em 5 etapas",
    [("1 · Coleta", "Notícias (5 portais, via API WordPress) + preços da PETR4 (B3/yfinance), com data e hora exatas."),
     ("2 · Sentimento (FinBERT-PT-BR)", "Cada manchete vira um número de sentimento; o dia é resumido no Índice de Sentimento (ISM)."),
     ("3 · Risco (GARCH)", "A série de retornos alimenta o GARCH(1,1), que estima a volatilidade condicional (o risco) do dia."),
     ("4 · Fusão de dados", "Sentimento + risco + retorno + categorias são concatenados num único vetor, todos defasados (t−1)."),
     ("5 · Classificação (XGBoost)", "O vetor entra no XGBoost, que projeta a direção do PRÓXIMO pregão; validação walk-forward."),
     ("Marcação temporal", "Notícia após as 17h conta para o pregão seguinte (Lead-Lag) — evita usar o futuro.")],
    depois_de="Arquitetura da Solução", sz_label=16.5, sz_desc=14.5, space=10)


# ============================================================================
#  9) FUSAO DE DADOS (slide 24) — reescrever a explicacao (mantendo a imagem)
# ============================================================================
for s in list(prs.slides):
    if H.title_text(s).startswith("Fusão de Dados"):
        for sh in list(s.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip().startswith("Os atributos heterogêneos"):
                tf = sh.text_frame; tf.clear()
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = ("Data Fusion = juntar informações de naturezas diferentes num ÚNICO vetor por pregão: "
                          "sentimento das notícias (FinBERT), risco (GARCH), retorno e categorias temáticas — "
                          "todos defasados em t−1. Esse vetor entra no classificador, que devolve UMA previsão de "
                          "direção para o pregão seguinte. A pergunta do estudo: esse vetor 'fundido' prevê melhor "
                          "do que usar apenas preços?")
                r.font.size = Pt(16)


# ============================================================================
#  8) PSEUDOCODIGO — 2 slides simples (apos a Separacao dos Dados)
# ============================================================================
add_algoritmo(
    "Algoritmo 1 — construção da base diária",
    require="Notícias N (2016–2026); Preços P da PETR4; Taxonomia T (7 cat., 152 termos)",
    ensure="Base diária D = {(x_t, y_t)}, pronta para o classificador",
    linhas=[
        (1, 0, "para cada notícia n ∈ N faça", ""),
        (2, 1, "c_n  ← categoria(n, T)", "taxonomia supervisionada"),
        (3, 1, "s_n  ← FinBERT-PT-BR(título_n)", "sentimento ∈ [−1, +1]"),
        (4, 1, "se hora(n) ≥ 17h: dia(n) ← próximo pregão", "Lead-Lag"),
        (5, 0, "fim para", ""),
        (6, 0, "para cada pregão t faça", ""),
        (7, 1, "ISM_t ← Σ(s_n · conf_n) / Σ conf_n", "índice diário de sentimento"),
        (8, 1, "R_t   ← ln(P_t / P_{t−1})", "log-retorno"),
        (9, 1, "σ_t   ← GARCH(1,1) sobre {R}", "volatilidade condicional"),
        (10, 1, "y_t   ← 1 se R_t > 0 senão 0", "rótulo: alta / baixa"),
        (11, 1, "x_t   ← [ R_{t−1}, σ_{t−1}, ISM_{t−1}, ISM^cat_{t−1} ]", "atributos em t−1"),
        (12, 0, "fim para", ""),
        (13, 0, "retornar D = {(x_t, y_t)}", ""),
    ],
    depois_de="Separação dos Dados")

add_algoritmo(
    "Algoritmo 2 — treino e avaliação walk-forward",
    require="Base D; nº de janelas W; grade de hiperparâmetros Θ",
    ensure="Métricas no teste (Acurácia, F1, AUC) e importância dos atributos",
    linhas=[
        (1, 0, "(D_tr, D_val, D_te) ← split cronológico 60/15/25", "sem embaralhar datas"),
        (2, 0, "para cada janela w = 1, …, W faça", "walk-forward"),
        (3, 1, "M_w   ← treinar XGBoost em D_tr^w", "minimiza a log-loss"),
        (4, 1, "θ*, δ ← selecionar em D_val^w", "hiperparâmetros e limiar"),
        (5, 0, "fim para", ""),
        (6, 0, "M ← modelo final com θ*", ""),
        (7, 0, "para cada pregão t ∈ D_te faça", ""),
        (8, 1, "p_t ← M(x_t)", "probabilidade de alta"),
        (9, 1, "ŷ_t ← 1 se p_t ≥ δ senão 0", ""),
        (10, 0, "fim para", ""),
        (11, 0, "avaliar ŷ vs y em D_te", "uma única vez: Acurácia, F1, AUC"),
        (12, 0, "comparar com o baseline (classe maj. ≈ 53%)", ""),
        (13, 0, "retornar métricas e importância dos atributos", ""),
    ],
    depois_de="Algoritmo 1 — construção")


# ============================================================================
#  10) SEPARACAO DOS DADOS (slide 27) — corrigir split com janelas reais
# ============================================================================
reformatar("Separação dos Dados", "Separação dos Dados (validação cronológica)",
    [("Por que cronológico (nunca aleatório)", "Prever o futuro com o passado. Embaralhar datas vazaria informação futura."),
     ("Treino — 60% · 1.566 pregões", "2016-01 → 2022-04. O modelo aprende os padrões."),
     ("Validação — 15% · 391 pregões", "2022-04 → 2023-11. Ajuste de hiperparâmetros e do limiar de decisão."),
     ("Teste — 25% · 653 pregões", "2023-11 → 2026-06. Consultado UMA única vez, ao final — o resultado honesto."),
     ("Walk-forward", "O processo se repete em janelas deslizantes, simulando o uso real dia após dia.")],
    sz_label=17, sz_desc=15, space=13)


# ============================================================================
#  10/11) METRICAS (slide 28) — acrescentar explicacao de AUC
# ============================================================================
reformatar("Métricas de Avaliação", "Métricas de Avaliação",
    [("Saída", "Direção do pregão seguinte (alta/baixa)."),
     ("Acurácia", "% de acertos. Só vale se SUPERAR a classe majoritária (baseline ≈ 53%)."),
     ("Precisão / Revocação", "Dos previstos como alta, quantos eram alta; e das altas reais, quantas capturamos."),
     ("F1-Score", "Média harmônica de precisão e revocação — desmascara viés de classe."),
     ("AUC-ROC", "Qualidade do RANKING (0–1): probabilidade de o modelo dar nota maior a uma alta real do que a uma baixa. 0,50 = acaso; quanto mais perto de 1, melhor."),
     ("Validação", "Split cronológico 60/15/25; teste de 653 pregões consultado uma única vez.")],
    sz_label=17, sz_desc=14.5, space=11)


# ============================================================================
#  11) RESULTADOS: DIRECAO (slide 31) — reescrever notas (AUC/ablacao/split)
# ============================================================================
for s in list(prs.slides):
    if H.title_text(s).startswith("Experimentos e Resultados"):
        for sh in list(s.shapes):
            if sh.has_text_frame and sh.text_frame.text.strip().startswith("Janela 2016-2026"):
                blocos = [
                 ("Split (janela 2016–2026)", "treino 2016–2022 (1.566) · validação 2022–2023 (391) · teste 2023–2026 (653)."),
                 ("Melhor modelo", "XGBoost com Data Fusion: acurácia 52,22% e AUC 0,514 (0,50 = acaso)."),
                 ("Ganho do sentimento", "+2,45 pp sobre 'apenas preços'; o sentimento é o atributo mais importante (0,344)."),
                 ("Leitura honesta", "O ganho direcional é MODESTO e não significativo (binomial p = 0,145) e não supera o baseline (~53%)."),
                 ("Ablação por categoria", "Remover uma categoria e medir a queda de acurácia. 'Sanções/Navegação' é a mais informativa (−3,68 pp)."),
                 ("† F1 inflado", "O SVM prevê quase sempre 'alta' — por isso reportamos várias métricas."),
                ]
                tf = sh.text_frame
                H.fill_definitions(tf, blocos, sz_label=14.5, sz_desc=13, space=7)


# ============================================================================
#  13/14) BASELINE — comparacao honesta (apos o slide 31)
# ============================================================================
novo_conteudo(
    "Comparação com o baseline",
    [("Baseline (classe majoritária)", "Chutar sempre 'alta' acerta ~53,1% (o mercado sobe na maioria dos dias)."),
     ("Direção — XGBoost Data Fusion", "52,22% (AUC 0,514): NÃO supera o baseline. Direção diária é intrinsecamente difícil (mercado eficiente)."),
     ("Volatilidade (Granger) — AQUI está o ganho", "Sentimento → volatilidade é altamente significativo em TODAS as defasagens (p < 0,001)."),
     ("Efeito assimétrico (regressão quantílica)", "Nos piores dias (τ=0,05), o sentimento eleva o retorno em +261 bps (p = 0,034) — viés de negatividade."),
     ("Mensagem central", "O sentimento antecipa a TURBULÊNCIA (risco) muito mais do que a direção — e esse é o resultado forte e defensável.")],
    depois_de="Experimentos e Resultados", sz_label=16, sz_desc=14, space=11)


# ============================================================================
#  15) DISCUSSAO DOS RESULTADOS (apos o slide 32)
# ============================================================================
novo_conteudo(
    "Discussão dos Resultados",
    [("Desafios observados", "Baixa relação sinal-ruído: ~85% das notícias não deslocam o preço de forma anormal."),
     ("Onde o sinal se concentra", "Nos ~15% de notícias que coincidem com rompimentos de estresse no GARCH (choque informacional)."),
     ("Direção", "Ganho modesto e não significativo — coerente com a hipótese de mercados eficientes."),
     ("Volatilidade e assimetria", "Resultados fortes e significativos: o sentimento antecipa risco e pesa mais nos dias ruins."),
     ("Robustez", "Séries com caudas pesadas (Jarque-Bera), estacionárias (ADF) e com efeito ARCH — o GARCH é justificado."),
     ("Implicação", "Para a PETR4, o sentimento é mais um radar de RISCO do que uma bússola de direção.")],
    depois_de="Modelagem Preditiva Direcional", sz_label=16, sz_desc=14, space=10)


# ============================================================================
#  16) ENRIQUECIMENTO (a partir da pagina "Sobre a Pesquisa"): WordPress,
#      analise de sentimento, GARCH (grafico), ablacao (tabela), disp. (grafico)
# ============================================================================
# 16a) Coleta com data e hora (WordPress REST API) — apos "Dados: Corpus de Noticias"
novo_conteudo(
    "Coleta com data e hora (WordPress REST API)",
    [("O acesso", "WordPress REST API (/wp-json/wp/v2/posts) dos 5 portais — histórico completo, sem a janela curta das APIs comerciais."),
     ("Marcação temporal precisa", "Cada notícia traz 'date' (horário de Brasília) e 'date_gmt' (UTC) — requisito crítico apontado pela banca."),
     ("Viabiliza o Lead-Lag", "Permite separar o que foi publicado ANTES e DEPOIS do fechamento (17h); o que sai após 17h conta para o pregão seguinte."),
     ("Cinco fontes", "InfoMoney, Exame, MoneyTimes, Petronotícias e Poder360 — com a fonte atribuída a cada notícia (evita dependência de fonte única).")],
    depois_de="Dados: Corpus de Notícias", sz_label=16.5, sz_desc=14.5, space=11)

# 16b) Analise de sentimento (com termos reais) — apos "Taxonomia Temática"
novo_conteudo(
    "Análise de sentimento — do texto ao número",
    [("O modelo", "FinBERT-PT-BR (BERT ajustado a finanças em português) — substitui um classificador genérico por um do domínio."),
     ("A leitura", "Cada manchete → rótulo (positivo / negativo / neutro) + confiança → índice em [−1, +1]."),
     ("Exemplo", "‘Petrobras anuncia dividendos recordes’ → Positivo · ‘Ataque a porto petroleiro reduz a oferta’ → negativo p/ o mercado, mas ALTA para a produtora."),
     ("Termos que guiam a relevância", "152 termos em 7 categorias — ex.: ‘Brent’, ‘OPEP’, ‘Estreito de Ormuz’, ‘dividendos’, ‘troca de presidente’, ‘embargo’."),
     ("Índice diário (ISM)", "As notícias do dia são agregadas (peso = confiança do modelo) num único número por pregão.")],
    depois_de="Taxonomia Temática", sz_label=16, sz_desc=14, space=8)

# 16c) GARCH (com grafico real) — apos "Notícias Divergentes"
add_imagem(
    "Volatilidade: o modelo GARCH(1,1)",
    "grafico_volatilidade_garch.png",
    depois_de="Notícias Divergentes",
    legenda="Volatilidade condicional estimada pelo GARCH(1,1) para a PETR4 (dados reais da pesquisa).",
    blocos_topo=[
        ("O que faz", "Estima a volatilidade condicional (o risco) e capta o agrupamento — dias turbulentos vêm em sequência."),
        ("Por que t-Student", "Os retornos têm caudas pesadas (Jarque-Bera rejeita a normalidade) e efeito ARCH (ARCH-LM, p<0,001) — o GARCH é justificado."),
    ])

# 16d) Ablacao por categoria (tabela real) — apos "Experimentos e Resultados"
add_tabela(
    "Ablação por categoria temática",
    "Remove-se UMA categoria de notícias, treina-se de novo e mede-se a QUEDA de acurácia. Modelo completo: 53,45%. Quanto maior a queda, mais aquela categoria importa.",
    ["Categoria removida", "Acurácia sem ela", "Impacto (pp)"],
    [("CAT5 — Acordos, Sanções e Navegação", "49,77%", "−3,68"),
     ("CAT6 — Liderança e Governança", "50,38%", "−3,06"),
     ("CAT7 — Macroeconomia e Energia", "50,69%", "−2,76"),
     ("CAT1 — Empresa e Ativo", "51,61%", "−1,84"),
     ("CAT2 — Mercado de Petróleo", "51,91%", "−1,53"),
     ("CAT3 — Geopolítica e Conflitos", "51,91%", "−1,53"),
     ("CAT4 — Oferta e Infraestrutura", "51,91%", "−1,53")],
    depois_de="Experimentos e Resultados",
    larguras=[6.0, 2.4, 2.1], destaque=1)

# 16e) Sentimento x Volatilidade (grafico real) — apos "Modelagem Preditiva Direcional"
add_imagem(
    "Sentimento × Volatilidade",
    "grafico_dispersao_sentimento_volatilidade.png",
    depois_de="Modelagem Preditiva Direcional",
    legenda="Relação entre o sentimento diário e a volatilidade — sustenta o achado de Granger (sentimento → volatilidade, p<0,001).")


# ============================================================================
#  Correcoes finais de ordenacao/consistencia
# ============================================================================
# (a) O divisor 'Introdução' deve vir ANTES das contextualizacoes (o substring
#     'Introdução' havia casado com a agenda). Move o divisor para logo apos a agenda.
for i, s in enumerate(prs.slides):
    if H.title_text(s) == "Introdução":
        move_after(i, "Estrutura da Apresentação")
        break

# (b) Reconciliar '25 publicações' -> 'estudos selecionados' (o protocolo diz 29).
for s in list(prs.slides):
    for sh in list(s.shapes):
        if sh.has_text_frame and "síntese das 25 publicações" in sh.text_frame.text:
            H.replace_in_shape(sh, "síntese das 25 publicações da RSL",
                               "síntese dos estudos selecionados na RSL")

prs.save(OUT)
print("Salvo:", OUT, "| total de slides:", len(list(prs.slides)))
