# -*- coding: utf-8 -*-
import copy
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.oxml.ns import qn
from pptx.parts.presentation import PresentationPart
from pptx.opc.packuri import PackURI

def _safe_next_slide_partname(self):
    usados = set()
    try:
        for p in self.package.iter_parts():
            usados.add(str(p.partname))
    except Exception:
        pass
    n = 1
    while ("/ppt/slides/slide%d.xml" % n) in usados:
        n += 1
    return PackURI("/ppt/slides/slide%d.xml" % n)
PresentationPart._next_slide_partname = property(_safe_next_slide_partname)

def shapes_with_text(slide):
    return [sh for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()]

def find_shape(slide, substr):
    for sh in slide.shapes:
        if sh.has_text_frame and substr in sh.text_frame.text:
            return sh
    return None

def slide_index_by_text(prs, substr):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and substr in sh.text_frame.text:
                return i
    return -1

def replace_in_shape(sh, old, new):
    for p in sh.text_frame.paragraphs:
        runtxt = "".join(r.text for r in p.runs)
        if old in runtxt:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new); return True
            novo = runtxt.replace(old, new)
            for i, r in enumerate(p.runs):
                r.text = novo if i == 0 else ""
            return True
    return False

def _sldIdLst(prs):
    return prs.slides._sldIdLst

def delete_slide(prs, index):
    lst = _sldIdLst(prs); ids = list(lst); sldId = ids[index]
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId); lst.remove(sldId)

def duplicate_slide(prs, index):
    src = prs.slides[index]; layout = src.slide_layout
    new = prs.slides.add_slide(layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))
    lst = _sldIdLst(prs); ids = list(lst); sld = ids[-1]
    lst.remove(sld); ids2 = list(lst)
    lst.insert(list(lst).index(ids2[index]), sld)
    return prs.slides[index + 1]

def reorder_by_keys(prs, key_func):
    lst = prs.slides._sldIdLst; ids = list(lst); slides = list(prs.slides)
    ordem = sorted(range(len(slides)), key=lambda i: key_func(i, slides[i]))
    for e in ids: lst.remove(e)
    for i in ordem: lst.append(ids[i])

def find_title(slide):
    from pptx.enum.shapes import PP_PLACEHOLDER
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE or ph.placeholder_format.idx == 0:
                return ph
        except Exception:
            pass
    txt = shapes_with_text(slide)
    txt.sort(key=lambda s: (s.top if s.top is not None else 0))
    return txt[0] if txt else None

def title_text(slide):
    t = find_title(slide)
    return t.text_frame.text.strip() if t else ""

def set_title_text(ts, novo):
    p0 = ts.text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = novo
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = novo

def fill_definitions(tf, blocos, sz_label=18, sz_desc=15, space=14):
    tf.clear(); tf.word_wrap = True
    for i, (label, desc) in enumerate(blocos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space); p.space_before = Pt(0)
        if label:
            r1 = p.add_run(); r1.text = label + "  "; r1.font.bold = True; r1.font.size = Pt(sz_label)
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(sz_desc)
        pPr = p._p.get_or_add_pPr()
        for tag in ('a:buChar', 'a:buAutoNum'):
            for e in pPr.findall(qn(tag)):
                pPr.remove(e)
        pPr.append(pPr.makeelement(qn('a:buNone'), {}))

def clear_body(slide):
    title = find_title(slide); tel = title._element
    for sh in list(slide.shapes):
        if sh._element is tel:
            continue
        if (sh.has_text_frame and sh.text_frame.text.strip()) or sh.has_table or sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)
    return title

def add_content_slide(prs, titulo, blocos, ref_text="Perguntas de Pesquisa",
                      left=0.9, top=1.7, width=11.4, height=5.0, **kw):
    ref = slide_index_by_text(prs, ref_text)
    slide = duplicate_slide(prs, ref)
    title = find_title(slide); tel = title._element
    for sh in list(slide.shapes):
        if sh._element is tel:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            sh._element.getparent().remove(sh._element)
    set_title_text(title, titulo)
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    fill_definitions(tb.text_frame, blocos, **kw)
    return slide, tb

def blank_slide(prs, titulo, ref_text="Perguntas de Pesquisa"):
    ref = slide_index_by_text(prs, ref_text)
    slide = duplicate_slide(prs, ref)
    title = clear_body(slide)
    set_title_text(title, titulo)
    return slide
